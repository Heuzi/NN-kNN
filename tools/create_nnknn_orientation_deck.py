from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
OUT_PATH = REPO_ROOT / "papers and presentations" / "NN-kNN repo orientation.pptx"


BG = RGBColor(250, 250, 247)
INK = RGBColor(31, 38, 45)
MUTED = RGBColor(96, 105, 112)
TEAL = RGBColor(22, 128, 119)
CORAL = RGBColor(207, 92, 73)
GOLD = RGBColor(214, 157, 61)
BLUE = RGBColor(70, 105, 177)
GREEN = RGBColor(72, 148, 100)
LIGHT_TEAL = RGBColor(220, 241, 237)
LIGHT_CORAL = RGBColor(249, 226, 221)
LIGHT_GOLD = RGBColor(249, 239, 211)
LIGHT_BLUE = RGBColor(225, 232, 248)
WHITE = RGBColor(255, 255, 255)


def set_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_textbox(slide, x, y, w, h, text="", size=24, color=INK, bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.55, 0.35, 12.1, 0.45, "NN-kNN", 13, TEAL, True)
    add_textbox(slide, 0.55, 0.82, 12.2, 0.82, title, 33, INK, True)
    if subtitle:
        add_textbox(slide, 0.58, 1.58, 11.7, 0.38, subtitle, 15, MUTED)


def add_footer(slide, idx):
    add_textbox(slide, 11.95, 7.1, 0.75, 0.22, f"{idx:02d}", 9, MUTED)


def add_bullets(slide, x, y, w, h, items, size=18, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(9)
    return box


def add_pill(slide, x, y, w, h, label, fill, font_color=INK):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = "Aptos"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = font_color
    return shape


def add_card(slide, x, y, w, h, heading, body, accent):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(225, 226, 222)
    add_textbox(slide, x + 0.22, y + 0.18, w - 0.44, 0.28, heading, 16, accent, True)
    add_bullets(slide, x + 0.22, y + 0.62, w - 0.44, h - 0.75, body, 12.5, INK)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=MUTED):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(2)
    line.line.end_arrowhead = True
    return line


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_textbox(slide, 0.72, 0.8, 11.6, 0.55, "NN-kNN", 22, TEAL, True)
    add_textbox(slide, 0.7, 1.48, 11.8, 1.1, "Neural Case-Based Learning", 42, INK, True)
    add_textbox(
        slide,
        0.75,
        2.62,
        10.8,
        0.62,
        "Repository orientation for classification, regression, and CartPole RL workflows",
        19,
        MUTED,
    )
    add_pill(slide, 0.78, 4.05, 2.55, 0.45, "learn retrieval", LIGHT_TEAL, TEAL)
    add_pill(slide, 3.55, 4.05, 2.95, 0.45, "weight features", LIGHT_GOLD, GOLD)
    add_pill(slide, 6.75, 4.05, 3.0, 0.45, "aggregate cases", LIGHT_BLUE, BLUE)
    add_pill(slide, 9.98, 4.05, 2.35, 0.45, "adapt output", LIGHT_CORAL, CORAL)
    add_footer(slide, 1)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "What NN-kNN Is Here", "A neural retrieval model that keeps case-based structure visible.")
    add_card(
        slide,
        0.75,
        2.25,
        3.75,
        3.65,
        "Case Base",
        ["Training examples remain stored as retrievable cases.", "Queries compare against all stored cases in representation space."],
        TEAL,
    )
    add_card(
        slide,
        4.8,
        2.25,
        3.75,
        3.65,
        "Glocal Weighting",
        ["Feature dimensions are reweighted during distance computation.", "Activations are normalized over cases with softmax or sparsemax."],
        GOLD,
    )
    add_card(
        slide,
        8.85,
        2.25,
        3.75,
        3.65,
        "Prediction",
        ["Classification sums activation into class probability mass.", "Regression can use NN-CDH after retrieval for adaptation."],
        BLUE,
    )
    add_footer(slide, 2)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "Retrieval Pipeline", "The maintained core stays common while task heads differ.")
    xs = [0.78, 2.9, 5.05, 7.25, 9.45]
    labels = ["Input", "Project", "Distance", "Normalize", "Aggregate"]
    colors = [LIGHT_TEAL, LIGHT_BLUE, LIGHT_GOLD, LIGHT_CORAL, RGBColor(227, 239, 225)]
    accents = [TEAL, BLUE, GOLD, CORAL, GREEN]
    for i, (x, label) in enumerate(zip(xs, labels)):
        add_pill(slide, x, 2.45, 1.55, 0.64, label, colors[i], accents[i])
        if i < len(xs) - 1:
            add_arrow(slide, x + 1.6, 2.77, xs[i + 1] - 0.1, 2.77)
    add_bullets(
        slide,
        1.05,
        4.25,
        11.2,
        1.45,
        [
            "`normalize_over_cases` selects case-level normalization behavior.",
            "Retrieval-only output and post-adaptation output are reported separately in regression work.",
            "The same core supports tabular classification, regression, and experimental RL actors/critics.",
        ],
        18,
    )
    add_footer(slide, 3)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "Maintained Entry Points", "Use workflow helpers rather than notebook-only implementations.")
    add_card(slide, 0.75, 2.0, 3.75, 3.9, "Core", ["model/nnknn_model.py", "model/nn_cdh.py", "datasets/reg_data.py"], TEAL)
    add_card(
        slide,
        4.8,
        2.0,
        3.75,
        3.9,
        "Supervised",
        ["model/regression_workflow.py", "model/classification_workflow.py", "tools/table1_nnknn_kfold.py"],
        BLUE,
    )
    add_card(
        slide,
        8.85,
        2.0,
        3.75,
        3.9,
        "Reinforcement Learning",
        ["model/rl_workflow.py", "model/nec_workflow.py", "model/nnknn_rl_workflow.py", "datasets/rl_tasks.py"],
        CORAL,
    )
    add_footer(slide, 4)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "Classification Workflow", "Current classification uses the shared IJCAI-26 core.")
    add_bullets(
        slide,
        0.85,
        2.1,
        5.75,
        3.9,
        [
            "Supported small datasets: iris, zebra, zebra_special, wine, breast_cancer, balance, digits.",
            "Supported image workflows: MNIST, CIFAR-10, SVHN.",
            "Baselines include kNN, MLP, ConvNet, and feature-kNN variants.",
            "Text tasks are deferred to keep dependencies out of this workflow.",
        ],
        18,
    )
    img = REPO_ROOT / "papers and presentations" / "random vs case bias.png"
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(7.0), Inches(2.05), width=Inches(5.2))
    add_footer(slide, 5)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "Regression Workflow", "Regression reporting centers on retrieval, adaptation, and repeated benchmarks.")
    add_card(
        slide,
        0.75,
        2.05,
        3.55,
        3.8,
        "Run",
        ["Set `task_type=\"regression\"` explicitly.", "Prefer `model/regression_workflow.py` for orchestration."],
        TEAL,
    )
    add_card(
        slide,
        4.9,
        2.05,
        3.55,
        3.8,
        "Compare",
        ["Baselines: kNN(X), Oracle kNN, MLKR+kNN, MLP.", "Keep the same split per method within a run."],
        BLUE,
    )
    add_card(
        slide,
        9.05,
        2.05,
        3.55,
        3.8,
        "Report",
        ["Separate retrieval-only and post-NN-CDH results.", "Use mean/std summaries from repeated runs."],
        GOLD,
    )
    add_footer(slide, 6)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "Table 1 Protocol", "The paper-style regression table is scripted and restartable.")
    add_bullets(
        slide,
        0.85,
        2.05,
        11.8,
        4.25,
        [
            "`tools/table1_nnknn_kfold.py` runs 5-fold CV across the current regression dataset list.",
            "Rows distinguish pure retrieval, adaptation, locality, and locality plus adaptation.",
            "Normalizers include softmax and sparsemax variants, with locality toggles.",
            "Write fresh timestamped folders under `results/`; keep failed folders for postmortem comparison.",
            "Recommended exports: summary_long.csv, runs_long.csv, table1_like.csv, transposed.csv, done.json.",
        ],
        19,
    )
    add_footer(slide, 7)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "RL Baseline Protocol", "CartPole comparisons use fixed environment-step budgets and periodic evaluation.")
    add_card(slide, 0.75, 2.0, 3.6, 3.8, "DQN", ["CleanRL-style baseline.", "Use best checkpoint plus end-of-budget evaluation."], BLUE)
    add_card(slide, 4.85, 2.0, 3.6, 3.8, "NEC", ["Repo-native NEC workflow.", "Current fast reference remains below success threshold."], TEAL)
    add_card(
        slide,
        8.95,
        2.0,
        3.6,
        3.8,
        "NN-kNN-RL",
        ["On-policy actor-critic with GAE.", "NN-kNN or MLP actor; MLP or NN-kNN value critic."],
        CORAL,
    )
    add_footer(slide, 8)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "Current RL Status", "Treat these artifacts as diagnostic, not solved benchmark claims.")
    add_bullets(
        slide,
        0.85,
        2.0,
        11.7,
        4.55,
        [
            "DQN fast run selected eval mean return: 110.85 at 110k steps; below the 475.0 success threshold.",
            "NEC fast run selected eval mean return: 450.55 at 150k steps; improved but still below threshold.",
            "NN-kNN-RL smoke path validates plumbing only; current NN-kNN critic comparison reached 369.5.",
            "Use `training_efficiency` before making sample-efficiency or paper-style claims.",
        ],
        20,
    )
    add_footer(slide, 9)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "Fast Validation", "Small checks keep routine development cheap.")
    add_bullets(
        slide,
        0.9,
        2.0,
        11.4,
        4.75,
        [
            "`python codex/smoke_test.py --mode imports`",
            "`python codex/smoke_test.py --mode train`",
            "`python codex/smoke_test.py --mode classification`",
            "`python codex/smoke_test.py --mode rl`",
            "`python codex/smoke_test.py --mode nec`",
            "`python codex/smoke_test.py --mode nnknn_rl`",
            "Use synthetic datasets for quick validation unless the task requires larger real datasets.",
        ],
        18,
    )
    add_footer(slide, 10)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "Working Rules", "A few guardrails prevent confusing experiments.")
    add_bullets(
        slide,
        0.85,
        2.0,
        11.8,
        4.7,
        [
            "Treat `Outdated...` notebooks as archival only.",
            "Use notebooks for inspection; put serious repeated reporting in workflow/helper Python files.",
            "Do not overwrite run folders; create new timestamped artifacts under `results/`.",
            "For RL, fixed budgets plus best-checkpoint selection are required for fair DQN, NEC, and NN-kNN-RL comparisons.",
            "If success thresholds are never reached, interpret the run as possible underfit or unsolved.",
        ],
        19,
    )
    add_footer(slide, 11)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "Recommended Next Slides", "Where to extend this deck for a talk or paper meeting.")
    add_bullets(
        slide,
        0.85,
        2.0,
        11.6,
        4.35,
        [
            "Add one architecture figure for `NN_KNN_Model` and `GlocalFeatureWeight`.",
            "Add a Table 1 result slide after the next full k-fold run.",
            "Add RL learning curves only after a solved or clearly plateaued fixed-budget run.",
            "Add qualitative retrieval examples for Iris/Zebra and one regression dataset.",
        ],
        20,
    )
    add_footer(slide, 12)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    return OUT_PATH


if __name__ == "__main__":
    path = build_deck()
    print(path)
